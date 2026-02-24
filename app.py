import os
import streamlit as st
from groq import Groq
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

# =============================
# CONFIGURATION
# =============================

DATA_FOLDER = "data"
MODEL_NAME = "llama-3.3-70b-versatile"

st.set_page_config(
    page_title="RAGenius-ChatBot",
    layout="wide"
)

# =============================
# CUSTOM CSS
# =============================

st.markdown("""
<style>

.stApp {
    background-color: #0F172A;
    color: white;
}

.sidebar .sidebar-content {
    background-color: #020617;
}

.chat-user {
    background-color: #1E293B;
    padding: 10px;
    border-radius: 10px;
    margin-bottom: 5px;
}

.chat-bot {
    background-color: #334155;
    padding: 10px;
    border-radius: 10px;
    margin-bottom: 10px;
}

</style>
""", unsafe_allow_html=True)

# =============================
# GROQ CLIENT
# =============================

groq_client = Groq(
    api_key=os.getenv("GROQ_API_KEY")
)

def call_llm(prompt):

    completion = groq_client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model=MODEL_NAME,
        temperature=0
    )

    return completion.choices[0].message.content


# =============================
# FILE LOADERS
# =============================

def load_pdf(path):
    reader = PdfReader(path)
    return "\n".join([p.extract_text() or "" for p in reader.pages])


def load_docx(path):
    doc = DocxDocument(path)
    return "\n".join([p.text for p in doc.paragraphs])


def load_ppt(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def load_txt(path):
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def load_documents():

    docs = []

    if not os.path.exists(DATA_FOLDER):
        return docs

    for file in os.listdir(DATA_FOLDER):

        path = os.path.join(DATA_FOLDER, file)

        try:

            if file.endswith(".pdf"):
                text = load_pdf(path)

            elif file.endswith(".docx"):
                text = load_docx(path)

            elif file.endswith(".pptx"):
                text = load_ppt(path)

            elif file.endswith(".txt"):
                text = load_txt(path)

            else:
                continue

            docs.append(Document(page_content=text))

        except:
            pass

    return docs


# =============================
# VECTOR STORE
# =============================

def create_vector_store():

    docs = load_documents()

    if not docs:
        return None

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=800,
        chunk_overlap=100
    )

    chunks = splitter.split_documents(docs)

    embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-MiniLM-L6-v2"
    )

    vectorstore = FAISS.from_documents(
        chunks,
        embeddings
    )

    return vectorstore


# =============================
# SESSION STATE INIT
# =============================

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None

if "chats" not in st.session_state:
    st.session_state.chats = {}

if "current_chat" not in st.session_state:
    st.session_state.current_chat = "Chat 1"
    st.session_state.chats["Chat 1"] = []

if "chat_counter" not in st.session_state:
    st.session_state.chat_counter = 1


# =============================
# QUERY FUNCTION
# =============================

def unified_query(question):

    vector_store = st.session_state.vector_store

    q_lower = question.lower()

    if "brief" in q_lower:
        instruction = "Give brief answer in 2 sentences."

    elif "explain" in q_lower:
        instruction = "Explain clearly in 5 sentences."

    else:
        instruction = "Answer clearly in 3 sentences."


    # NORMAL MODE
    if vector_store is None:

        prompt = f"""
You are an AI assistant.

{instruction}

Question:
{question}

Answer:
"""

        return call_llm(prompt)


    # RAG MODE
    docs = vector_store.similarity_search(question, k=3)

    if not docs:
        return "No details found in documents."

    context = "\n".join([d.page_content for d in docs])

    prompt = f"""
Use ONLY context.

{instruction}

Context:
{context}

Question:
{question}

Answer:
"""

    return call_llm(prompt)


# =============================
# SIDEBAR
# =============================

st.sidebar.title("💬 Chats")

# NEW CHAT
if st.sidebar.button("➕ New Chat"):

    st.session_state.chat_counter += 1

    chat_name = f"Chat {st.session_state.chat_counter}"

    st.session_state.chats[chat_name] = []

    st.session_state.current_chat = chat_name

    st.rerun()


# SEARCH
search = st.sidebar.text_input("Search chat")


# CHAT LIST
for chat in st.session_state.chats:

    if search.lower() in chat.lower():

        if st.sidebar.button(chat):
            st.session_state.current_chat = chat
            st.rerun()


# =============================
# MAIN UI
# =============================

st.title("🧠 RAGenius-ChatBot")

st.subheader(st.session_state.current_chat)

chat_area = st.container()

with chat_area:

    for msg in st.session_state.chats[st.session_state.current_chat]:

        if msg["role"] == "user":

            st.markdown(
                f'<div class="chat-user">🧑 {msg["content"]}</div>',
                unsafe_allow_html=True
            )

        else:

            st.markdown(
                f'<div class="chat-bot">🤖 {msg["content"]}</div>',
                unsafe_allow_html=True
            )


# =============================
# INPUT AREA
# =============================

col1, col2 = st.columns([4,1])

with col1:

    with st.form("input", clear_on_submit=True):

        question = st.text_input("Ask question")

        send = st.form_submit_button("Send")


with col2:

    uploaded = st.file_uploader(
        "Upload file",
        type=["pdf","docx","pptx","txt"]
    )


# =============================
# FILE UPLOAD
# =============================

if uploaded:

    os.makedirs(DATA_FOLDER, exist_ok=True)

    path = os.path.join(DATA_FOLDER, uploaded.name)

    with open(path, "wb") as f:
        f.write(uploaded.getbuffer())

    st.success("Uploaded")

    st.session_state.vector_store = create_vector_store()

    st.success("Knowledge updated")


# =============================
# MESSAGE PROCESS
# =============================

if send and question:

    current = st.session_state.current_chat

    st.session_state.chats[current].append(
        {"role":"user","content":question}
    )

    with st.spinner("Thinking..."):

        answer = unified_query(question)

    st.session_state.chats[current].append(
        {"role":"assistant","content":answer}
    )

    st.rerun()